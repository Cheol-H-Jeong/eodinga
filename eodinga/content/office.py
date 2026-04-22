from __future__ import annotations

from io import BytesIO
from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from eodinga.content.base import ParsedContent, ParserSpec, make_parsed_content
from eodinga.core.fs import read_bytes


def parse_docx(path: Path, max_body_chars: int) -> ParsedContent:
    document = Document(BytesIO(read_bytes(path)))
    paragraphs = [
        paragraph.text.strip()
        for paragraph in document.paragraphs
        if paragraph.text.strip()
    ]
    body_text = "\n".join(paragraphs)
    head_text = "\n".join(paragraphs[:5]).strip()
    title = paragraphs[0][:200] if paragraphs else path.stem
    return make_parsed_content(
        title=title,
        head_text=head_text,
        body_text=body_text,
        max_body_chars=max_body_chars,
    )


def parse_pptx(path: Path, max_body_chars: int) -> ParsedContent:
    presentation = Presentation(BytesIO(read_bytes(path)))
    texts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", "").strip()
            if text:
                texts.append(text)
    body_text = "\n".join(texts)
    head_text = "\n".join(texts[:5]).strip()
    title = texts[0][:200] if texts else path.stem
    return make_parsed_content(
        title=title,
        head_text=head_text,
        body_text=body_text,
        max_body_chars=max_body_chars,
    )


def parse_xlsx(path: Path, max_body_chars: int) -> ParsedContent:
    workbook = load_workbook(BytesIO(read_bytes(path)), read_only=True, data_only=True)
    rows: list[str] = []
    for sheet in workbook.worksheets[:5]:
        rows.append(f"[{sheet.title}]")
        for row in sheet.iter_rows(min_row=1, max_row=100, values_only=True):
            values = [
                str(value).strip()
                for value in row
                if value is not None and str(value).strip()
            ]
            if values:
                rows.append(" | ".join(values))
    body_text = "\n".join(rows)
    head_text = "\n".join(rows[:5]).strip()
    title = rows[0][:200] if rows else path.stem
    return make_parsed_content(
        title=title,
        head_text=head_text,
        body_text=body_text,
        max_body_chars=max_body_chars,
    )


def get_docx_parser_spec() -> ParserSpec:
    return ParserSpec(name="docx", extensions=frozenset({"docx"}), parse=parse_docx)


def get_pptx_parser_spec() -> ParserSpec:
    return ParserSpec(name="pptx", extensions=frozenset({"pptx"}), parse=parse_pptx)


def get_xlsx_parser_spec() -> ParserSpec:
    return ParserSpec(name="xlsx", extensions=frozenset({"xlsx"}), parse=parse_xlsx)
